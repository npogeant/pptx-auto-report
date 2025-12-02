from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33) # 16:9 width
    prs.slide_height = Inches(7.5)  # 16:9 height
    return prs

def set_font(shape, size=32, bold=False):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Century Gothic"
            run.font.size = size
            run.font.bold = bold
            run.font.color.rgb = RGBColor(0, 0, 0)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    title_shape = slide.shapes.title
    title_shape.text = title
    set_font(title_shape, size=Pt(40), bold=True)
    
    # Align the title shape centered horizontally on slide
    title_shape.width = Inches(12.33)
    title_shape.left = (prs.slide_width - title_shape.width) // 2
    title_shape.top = Inches(2)
    title_shape.height = Pt(60)  # Height based on 40pt font size with padding

    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    set_font(subtitle_shape, size=Pt(24), bold=False)
    
    # Align the subtitle shape centered horizontally on slide
    subtitle_shape.width = Inches(12.33)
    subtitle_shape.left = (prs.slide_width - subtitle_shape.width) // 2
    subtitle_shape.top = Inches(4)
    subtitle_shape.height = Pt(36)  # Height based on 24pt font size with padding

    return slide

def add_chart_slide(prs, chart_data, chart_type=XL_CHART_TYPE.LINE, title=""):
    slide_layout = prs.slide_layouts[5]  # blank slide
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    set_font(title_shape, size=Pt(30), bold=True)
    
    # Align the title shape centered horizontally on slide
    title_shape.width = Inches(12.33)
    title_shape.left = (prs.slide_width - title_shape.width) // 2
    title_shape.top = Inches(0.8)
    title_shape.height = Pt(45)  # Height based on 40pt font size with padding

    # Layout: chart on the left, two cards on the right
    left_chart = Inches(1)
    chart_width = Inches(8.5)
    chart_height = Inches(5)
    top = Inches(2)

    chart = slide.shapes.add_chart(
        chart_type,
        left_chart, top, chart_width, chart_height,
        chart_data
    ).chart

    # ---- Remove title ----
    chart.has_title = False
    
    # ---- Legend above chart ----
    chart.has_legend = False

    # ---- Plot area ----
    plot = chart.plots[0]

    # Enable smoothing for the line chart and hide plot-level data-labels
    plot.has_data_labels = False

    # Define colors for each series (left-to-right)
    series_colors = [
        RGBColor(46, 117, 182),   # Blue
        RGBColor(237, 125, 49)    # Orange
    ]

    # Compute evolutions (delta, pct) for each series.
    # Prefer chart_data.evolution (set by chart_builder), otherwise try to compute from chart_data.series
    evolutions = []
    if hasattr(chart_data, 'evolution') and isinstance(chart_data.evolution, (list, tuple)):
        evolutions = list(chart_data.evolution)
    else:
        try:
            if hasattr(chart_data, 'series'):
                for s in chart_data.series:
                    vals = getattr(s, 'values', None)
                    if vals is None:
                        vals = list(s)
                    numeric = [v for v in vals if isinstance(v, (int, float))]
                    if len(numeric) >= 1:
                        start = float(numeric[0])
                        end = float(numeric[-1])
                        delta = end - start
                        pct = (delta / start) * 100.0 if start != 0 else None
                        evolutions.append((float(delta), None if pct is None else float(pct)))
                    else:
                        evolutions.append((None, None))
        except Exception:
            evolutions = []

    # Apply styling to all series
    for idx, series in enumerate(plot.series):
        color = series_colors[idx % len(series_colors)]

        # Line thickness
        series.format.line.width = Pt(2.5)

        # Apply smoothing to the line (lissage)
        try:
            series.smooth = True
        except Exception:
            pass

        # Marker style
        series.marker.style = XL_MARKER_STYLE.CIRCLE
        series.marker.size = 6

        # Line color
        series.format.line.color.rgb = color

        # Marker color
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = color
        # Remove marker outline (contour)
        try:
            series.marker.format.line.width = Pt(0)
            # try to clear the line fill if supported
            try:
                series.marker.format.line.fill.background()
            except Exception:
                pass
        except Exception:
            pass

        # ---- Data labels ----
        series.has_data_labels = True
        data_labels = series.data_labels
        data_labels.font.name = "Century Gothic"
        data_labels.font.size = Pt(10)
        data_labels.font.color.rgb = color

        # place data labels above markers and show values
        try:
            data_labels.position = XL_LABEL_POSITION.ABOVE
            data_labels.show_value = True
        except Exception:
            pass

    # ---- Add right-side cards showing totals for each series ----
    right_margin = Inches(0.5)
    cards_left = left_chart + chart_width + right_margin
    cards_width = prs.slide_width - cards_left - Inches(0.5)
    card_height = Inches(2.0)
    card_gap = Inches(0.3)

    # colors for positive/negative evolution
    green = RGBColor(0, 153, 76)
    red = RGBColor(204, 0, 0)

    for i in range(2):
        card_color = series_colors[i % len(series_colors)]

        card_top = top + i * (card_height + card_gap)
        card = slide.shapes.add_textbox(cards_left, card_top, cards_width, card_height)
        tf = card.text_frame
        tf.clear()

        # Big value -> show percentage only (in place of absolute delta)
        p_val = tf.paragraphs[0]
        p_val.text = ''
        run = p_val.add_run()
        # get evolution (delta, pct)
        delta = None
        pct = None
        if i < len(evolutions):
            ev = evolutions[i]
            if isinstance(ev, (list, tuple)) and len(ev) >= 2:
                delta, pct = ev[0], ev[1]

        if pct is None:
            display_value = "-"
        else:
            try:
                sign = '+' if pct >= 0 else '-'
                display_value = f"{sign}{abs(pct):.0f}%"
            except Exception:
                display_value = str(pct)

        run.text = display_value
        run.font.name = "Century Gothic"
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = card_color

        # Label (series name)
        p_label = tf.add_paragraph()
        # get series name from chart_data if possible, else from chart
        series_name = None
        if hasattr(chart_data, 'series_names') and i < len(chart_data.series_names):
            series_name = chart_data.series_names[i]
        else:
            try:
                if hasattr(chart_data, 'series') and i < len(chart_data.series):
                    series_name = getattr(chart_data.series[i], 'name', None)
            except Exception:
                series_name = None

        if not series_name:
            try:
                series_name = plot.series[i].name
            except Exception:
                series_name = f"Series {i+1}"

        p_label.text = series_name
        p_label.font.name = "Century Gothic"
        p_label.font.size = Pt(14)
        p_label.font.bold = False
        p_label.font.color.rgb = card_color

    # ---- Remove gridlines ----
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = False

    # ---- Hide Y axis ----
    value_axis.visible = False

    # ---- Category axis font ----
    category_axis = chart.category_axis
    category_axis.tick_labels.font.name = "Century Gothic"
    category_axis.tick_labels.font.size = Pt(12)

    # Remove X axis line (contour) while keeping tick labels
    try:
        category_axis.format.line.width = Pt(0)
        try:
            category_axis.format.line.fill.background()
        except Exception:
            pass
    except Exception:
        try:
            # older/newer variants
            category_axis.line_format.width = Pt(0)
        except Exception:
            pass

    return slide

def save_presentation(prs, output_path="../output/report.pptx"):
    prs.save(output_path)
